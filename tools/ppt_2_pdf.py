import os
import tempfile
import io
import math
from typing import Any, Dict, List, Optional, Union

from dify_plugin import Tool
from dify_plugin.file.file import File

# 导入依赖库
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX, MSO_LINE_DASH_STYLE
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Table, TableStyle, Paragraph, Frame, KeepInFrame
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    DEPENDENCIES_AVAILABLE = True
except ImportError:
    DEPENDENCIES_AVAILABLE = False

# 常量定义
EMU_PER_INCH = 914400
PT_PER_INCH = 72
EMU_TO_PT = PT_PER_INCH / EMU_PER_INCH

class PptToPdfTool(Tool):
    """
    Advanced PPT to PDF Converter (Pure Python) V4.
    Fixes: 'MSO_SHAPE_TYPE' attribute error, Connectors logic, and Layouts.
    """

    def _invoke(self, tool_parameters: dict[str, Any]) -> Any:
        if not DEPENDENCIES_AVAILABLE:
            yield self.create_text_message("Error: Required libraries (python-pptx, reportlab, Pillow) are missing.")
            return

        input_file = tool_parameters.get("input_file")
        if not input_file:
            yield self.create_text_message("Error: Input file is required.")
            return

        if not input_file.extension or input_file.extension.lower() not in ['.pptx']:
            yield self.create_text_message("Error: Only .pptx files are supported.")
            return

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                input_path = os.path.join(temp_dir, input_file.filename)
                with open(input_path, "wb") as f:
                    f.write(input_file.blob)
                
                output_filename = os.path.splitext(input_file.filename)[0] + ".pdf"
                output_path = os.path.join(temp_dir, output_filename)

                converter = PptPdfEngine(input_path, output_path)
                result = converter.convert()

                if not result["success"]:
                    yield self.create_text_message(f"Conversion Failed: {result['message']}")
                    return

                with open(output_path, 'rb') as f:
                    pdf_content = f.read()

                yield self.create_text_message("PPT conversion successful.")
                yield self.create_blob_message(
                    blob=pdf_content,
                    meta={
                        "filename": output_filename,
                        "mime_type": "application/pdf"
                    }
                )

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield self.create_text_message(f"System Error: {str(e)}")

class PptPdfEngine:
    def __init__(self, input_path: str, output_path: str):
        self.input_path = input_path
        self.output_path = output_path
        self.font_name = "CustomChineseFont"
        self.font_bold_name = "CustomChineseFont"
        self._register_fonts()

    def _register_fonts(self):
        """字体注册逻辑"""
        try:
            current_dir = os.path.dirname(os.path.abspath(__file__))
            # 假设字体在 ../fonts/
            font_path = os.path.join(os.path.dirname(current_dir), "fonts", "chinese_font.ttc")
            
            if not os.path.exists(font_path):
                font_path = os.path.join(current_dir, "fonts", "chinese_font.ttc")

            if os.path.exists(font_path):
                pdfmetrics.registerFont(TTFont(self.font_name, font_path))
                # 简单复用作为粗体
                pdfmetrics.registerFont(TTFont(self.font_bold_name, font_path)) 
            else:
                self.font_name = "Helvetica"
                self.font_bold_name = "Helvetica-Bold"
        except Exception:
            self.font_name = "Helvetica"
            self.font_bold_name = "Helvetica-Bold"

    def convert(self) -> Dict[str, Any]:
        try:
            prs = Presentation(self.input_path)
            slide_width_pt = prs.slide_width * EMU_TO_PT
            slide_height_pt = prs.slide_height * EMU_TO_PT
            
            c = canvas.Canvas(self.output_path, pagesize=(slide_width_pt, slide_height_pt))
            
            for slide in prs.slides:
                self._process_slide(c, slide, slide_height_pt)
                c.showPage()
            
            c.save()
            return {"success": True, "message": "OK"}
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"success": False, "message": str(e)}

    def _process_slide(self, c: canvas.Canvas, slide: Any, page_height: float):
        # 1. 背景绘制 (简化处理)
        try:
            bg = slide.background
            if bg and bg.fill and bg.fill.type: 
                color = self._get_solid_fill_color(bg.fill)
                if color:
                    c.setFillColor(color)
                    c.rect(0, 0, c._pagesize[0], c._pagesize[1], fill=1, stroke=0)
        except:
            pass

        # 2. 形状绘制
        if slide.shapes:
            for shape in slide.shapes:
                self._render_shape_recursive(c, shape, 0, 0, page_height)

    def _render_shape_recursive(self, c: canvas.Canvas, shape: Any, x_offset: float, y_offset: float, page_height: float):
        if hasattr(shape, 'visible') and not shape.visible:
            return

        # 坐标计算 (EMU)
        try:
            current_x_emu = x_offset + shape.left
            current_y_emu = y_offset + shape.top
        except AttributeError:
            current_x_emu = x_offset
            current_y_emu = y_offset

        # --- 分类处理 ---

        # 1. 组合 (Group)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                self._render_shape_recursive(c, sub_shape, current_x_emu, current_y_emu, page_height)
            return

        # 2. 连接线/线条 (Line/Connector) - 【修复点】使用 hasattr 检测
        # 只要有 begin_x 和 end_x 属性，我们就认为它是一条线
        if hasattr(shape, 'begin_x') and hasattr(shape, 'end_x'):
           self._draw_connector(c, shape, x_offset, y_offset, page_height)
           return

        # 3. 常规形状 (AutoShape, TextBox, Picture, Table)
        try:
            x = current_x_emu * EMU_TO_PT
            w = shape.width * EMU_TO_PT
            h = shape.height * EMU_TO_PT
            y = page_height - (current_y_emu * EMU_TO_PT) - h 
        except AttributeError:
            return

        # 3.1 形状背景与边框 (修复背景丢失)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
             self._draw_shape_background(c, shape, x, y, w, h)

        # 3.2 文本内容
        if shape.has_text_frame and shape.text_frame.text.strip():
            self._draw_smart_text_box(c, shape, x, y, w, h)
        
        # 3.3 表格
        elif shape.has_table:
            self._draw_exact_table(c, shape.table, x, y, w, h, page_height)

        # 3.4 图片
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_blob = shape.image.blob
                img_reader = ImageReader(io.BytesIO(image_blob))
                c.drawImage(img_reader, x, y, width=w, height=h, mask='auto', preserveAspectRatio=True)
            except Exception:
                pass

    def _draw_connector(self, c: canvas.Canvas, shape: Any, x_offset: float, y_offset: float, page_height: float):
        """绘制线条"""
        try:
            # begin_x/y 属性通常是相对于 Slide 的绝对坐标 (但如果 API 版本变动可能不同)
            # 受限于 python-pptx 的实现，我们暂时忽略 Group 对连接线的偏移，因为通常连接线使用绝对坐标
            bx = shape.begin_x * EMU_TO_PT
            by = page_height - (shape.begin_y * EMU_TO_PT)
            ex = shape.end_x * EMU_TO_PT
            ey = page_height - (shape.end_y * EMU_TO_PT)

            c.saveState()
            
            line_width = 1
            line_color = colors.black
            
            if hasattr(shape, 'line'):
                # 填充色
                if shape.line.fill and shape.line.fill.type:
                     lc = self._get_solid_fill_color(shape.line.fill)
                     if lc: line_color = lc
                # 宽度
                if shape.line.width:
                    line_width = shape.line.width.pt
                # 虚线
                self._apply_dash_style(c, shape.line)

            c.setStrokeColor(line_color)
            c.setLineWidth(line_width)
            c.line(bx, by, ex, ey)
            
            c.restoreState()
        except AttributeError:
            pass

    def _draw_shape_background(self, c: canvas.Canvas, shape: Any, x, y, w, h):
        """绘制背景和边框"""
        fill_color = None
        line_color = None
        line_width = 0

        c.saveState()

        # 1. 填充
        if hasattr(shape, 'fill'):
            fill_color = self._get_solid_fill_color(shape.fill)

        # 2. 边框
        if hasattr(shape, 'line'):
             if shape.line.fill and shape.line.fill.type:
                 line_color = self._get_solid_fill_color(shape.line.fill)
             
             if hasattr(shape.line, 'width') and shape.line.width:
                 line_width = shape.line.width.pt
             
             if line_color and line_width > 0:
                 self._apply_dash_style(c, shape.line)

        # 设置绘图属性
        if fill_color:
            c.setFillColor(fill_color)
        else:
            c.setFillColor(colors.Color(0,0,0,alpha=0))

        if line_color and line_width > 0:
            c.setStrokeColor(line_color)
            c.setLineWidth(line_width)
        else:
            c.setStrokeColor(colors.Color(0,0,0,alpha=0))

        # 绘制矩形
        c.rect(x, y, w, h, fill=1 if fill_color else 0, stroke=1 if line_color else 0)
        
        c.restoreState()

    def _apply_dash_style(self, c: canvas.Canvas, line_fmt: Any):
        """Mapping PPT dash styles to PDF"""
        try:
            if not hasattr(line_fmt, 'dash_style') or not line_fmt.dash_style:
                return

            style = line_fmt.dash_style
            width = line_fmt.width.pt if (line_fmt and hasattr(line_fmt, 'width') and line_fmt.width) else 1
            
            # MSO_LINE_DASH_STYLE 键值映射
            dash_map = {
                MSO_LINE_DASH_STYLE.DASH: [width * 4, width * 3],
                MSO_LINE_DASH_STYLE.DASH_DOT: [width * 4, width * 3, width * 1, width * 3],
                MSO_LINE_DASH_STYLE.DOT: [width * 1, width * 3],
                MSO_LINE_DASH_STYLE.LONG_DASH: [width * 8, width * 3],
                MSO_LINE_DASH_STYLE.LONG_DASH_DOT: [width * 8, width * 3, width * 1, width * 3],
                MSO_LINE_DASH_STYLE.ROUND_DOT: [1, width * 4],
                MSO_LINE_DASH_STYLE.SQUARE_DOT: [width * 1, width * 1],
                MSO_LINE_DASH_STYLE.SOLID: []
            }
            
            pattern = dash_map.get(style)
            if pattern is not None: # pattern might be empty list [] for solid
                c.setDash(pattern)
                
        except Exception:
            pass

    def _draw_smart_text_box(self, c: canvas.Canvas, shape: Any, x, y, w, h):
        text_frame = shape.text_frame
        styles = getSampleStyleSheet()
        flowables = []
        
        for paragraph in text_frame.paragraphs:
            if not paragraph.text and not paragraph.runs:
                flowables.append(Paragraph("<br/>", styles["Normal"]))
                continue

            font_size = 10 
            font_color = colors.black
            is_bold = False
            
            if paragraph.runs:
                run = paragraph.runs[0]
                if run.font.size: font_size = run.font.size.pt
                if run.font.bold: is_bold = True
                
                if run.font.color:
                     c_val = self._get_color_from_font(run.font)
                     if c_val: font_color = c_val

            used_font = self.font_bold_name if is_bold else self.font_name

            style = ParagraphStyle(
                name=f'P_{id(paragraph)}',
                parent=styles['Normal'],
                fontName=used_font,
                fontSize=font_size,
                textColor=font_color,
                leading=font_size * 1.2,
                wordWrap='CJK',
                alignment=self._map_alignment(paragraph.alignment)
            )
            
            txt = paragraph.text if paragraph.text else ""
            txt = txt.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
            
            flowables.append(Paragraph(txt, style))

        if not flowables:
            return

        m_l = shape.text_frame.margin_left * EMU_TO_PT if hasattr(shape.text_frame, 'margin_left') else 5
        m_r = shape.text_frame.margin_right * EMU_TO_PT if hasattr(shape.text_frame, 'margin_right') else 5
        m_t = shape.text_frame.margin_top * EMU_TO_PT if hasattr(shape.text_frame, 'margin_top') else 5
        m_b = shape.text_frame.margin_bottom * EMU_TO_PT if hasattr(shape.text_frame, 'margin_bottom') else 5

        # 容错：如果margin计算后宽度不足，强制给点空间
        draw_w = max(10, w - m_l - m_r)
        draw_h = max(10, h - m_t - m_b)

        frame = Frame(x, y, w, h, showBoundary=0, topPadding=m_t, leftPadding=m_l, rightPadding=m_r, bottomPadding=m_b)
        story = [KeepInFrame(draw_w, draw_h, flowables, mode='shrink')]
        frame.addFromList(story, c)

    def _draw_exact_table(self, c: canvas.Canvas, ppt_table: Any, x, y, w, h, page_height):
        if not ppt_table.rows: return

        data = []
        row_heights = []
        for row in ppt_table.rows:
            row_data = []
            row_heights.append(row.height * EMU_TO_PT)
            for cell in row.cells:
                txt = cell.text_frame.text.strip() if cell.text_frame else ""
                row_data.append(txt)
            data.append(row_data)

        col_widths = [col.width * EMU_TO_PT for col in ppt_table.columns]
        
        processed_data = []
        base_style = ParagraphStyle(name='TB', fontName=self.font_name, fontSize=9, leading=11, wordWrap='CJK')
        
        for row in data:
            new_row = []
            for txt in row:
                safe_txt = txt.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                new_row.append(Paragraph(safe_txt, base_style))
            processed_data.append(new_row)

        rl_table = Table(processed_data, colWidths=col_widths, rowHeights=row_heights)
        rl_table.setStyle(TableStyle([
            ('FONTNAME', (0, 0), (-1, -1), self.font_name),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]))

        t_w, t_h = rl_table.wrap(w, h)
        top_y = y + h
        draw_y = top_y - t_h
        rl_table.drawOn(c, x, draw_y)

    def _get_solid_fill_color(self, fill_obj):
        try:
            if hasattr(fill_obj, 'fore_color'):
                 if fill_obj.fore_color.type == MSO_COLOR_TYPE.RGB:
                     return self._rgb_to_color(fill_obj.fore_color.rgb)
                 elif fill_obj.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                     # 对于主题色，简单映射为浅灰色背景，以便能看到边框
                     return colors.Color(0.9, 0.9, 0.9) 
        except Exception:
            pass
        return None
        
    def _get_color_from_font(self, font_obj):
        try:
            if font_obj.color.type == MSO_COLOR_TYPE.RGB:
                return self._rgb_to_color(font_obj.color.rgb)
        except Exception:
            pass
        return None

    def _rgb_to_color(self, rgb):
        return colors.Color(rgb[0]/255.0, rgb[1]/255.0, rgb[2]/255.0)

    def _map_alignment(self, ppt_align):
        if ppt_align == 1: return TA_CENTER
        if ppt_align == 2: return TA_RIGHT
        if ppt_align == 3: return TA_JUSTIFY
        return TA_LEFT